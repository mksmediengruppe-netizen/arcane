"""
ARCANE Tool Executor
Dispatches tool calls to the appropriate worker/handler.
Manages sandbox sessions, file operations, and browser state.
"""

from __future__ import annotations

import asyncio
import json
import os
import subprocess
import time
import tempfile
from typing import Any, Optional

from core.tool_registry import ToolRegistry
from core.sandbox import execute_sandboxed, execute as execute_unified, is_command_safe
from shared.utils.error_analyzer import analyze_error
from shared.utils.logger import get_logger, log_with_data

logger = get_logger("core.tool_executor")

# ═══════════════════════════════════════════════════════════════════════════════
# PATH SECURITY — prevent path traversal attacks
# ═══════════════════════════════════════════════════════════════════════════════

ALLOWED_BASE_DIRS = [
    "/root/arcane",
    "/root/workspace",
    "/tmp",
    "/home/arcane_sandbox",
    "/usr/lib",
    "/usr/local",
    # "/home/ubuntu/projects",  # P4-FIX BUG-005: removed, use /root/workspace/ instead
]


def _validate_file_path(path: str) -> tuple[bool, str]:
    """Validate that a file path is within allowed directories.

    Uses os.path.commonpath() to prevent path traversal attacks.
    startswith() is vulnerable: '/tmpx/evil.txt'.startswith('/tmp') == True.
    commonpath() correctly checks that paths are truly nested.

    Returns (is_valid, resolved_path_or_error_message).
    """
    if not path:
        return False, "Path is empty"

    real_path = os.path.realpath(os.path.expanduser(path))

    for allowed in ALLOWED_BASE_DIRS:
        try:
            common = os.path.commonpath([real_path, allowed])
            if common == allowed:
                return True, real_path
        except ValueError:
            continue

    return False, f"Access denied: {path} is outside allowed directories ({', '.join(ALLOWED_BASE_DIRS)})"


class SandboxSession:
    """Manages an isolated shell session for command execution."""

    def __init__(self, session_id: str = "default", working_dir: str = "/root/workspace"):
        self.session_id = session_id
        self.working_dir = working_dir
        self._history: list[dict] = []

    async def execute(self, command: str, timeout: int = 30, working_dir: str = None) -> dict:
        """Execute a command in the sandbox with optional isolation."""
        cwd = working_dir or self.working_dir
        os.makedirs(cwd, exist_ok=True)

        # Use unified execution — Docker first, fallback to su-based
        result = await execute_unified(
            command=command,
            working_dir=cwd,
            session_id=self.session_id,
            timeout=timeout,
        )

        self._history.append({"command": command, **result})
        return result


class ToolExecutor:
    """
    Executes tool calls from the agent loop.
    Routes each tool to its appropriate handler.
    """

    def __init__(self, registry: ToolRegistry, project_dir: str = "/root/workspace"):
        self._registry = registry
        self._project_dir = project_dir
        self._sessions: dict[str, SandboxSession] = {}
        self._browser = None  # Lazy-initialized BrowserWorker
        self._setup_builtin_handlers()

    def _setup_builtin_handlers(self):
        """Register built-in tool handlers."""
        self._registry.register_handler("shell_exec", self._handle_shell_exec)
        self._registry.register_handler("ssh_exec", self._handle_ssh_exec)
        self._registry.register_handler("shell_view", self._handle_shell_view)
        self._registry.register_handler("file_read", self._handle_file_read)
        self._registry.register_handler("file_write", self._handle_file_write)
        self._registry.register_handler("file_edit", self._handle_file_edit)
        self._registry.register_handler("file_append", self._handle_file_append)
        self._registry.register_handler("glob", self._handle_glob)
        self._registry.register_handler("grep", self._handle_grep)
        self._registry.register_handler("message", self._handle_message)
        self._registry.register_handler("plan", self._handle_plan)
        # Browser tools
        self._registry.register_handler("browser_navigate", self._handle_browser_navigate)
        self._registry.register_handler("browser_view", self._handle_browser_view)
        self._registry.register_handler("browser_click", self._handle_browser_click)
        self._registry.register_handler("browser_input", self._handle_browser_input)
        self._registry.register_handler("browser_scroll", self._handle_browser_scroll)
        self._registry.register_handler("browser_select", self._handle_browser_select)
        self._registry.register_handler("browser_find", self._handle_browser_find)
        self._registry.register_handler("browser_save_image", self._handle_browser_save_image)
        self._registry.register_handler("browser_press_key", self._handle_browser_press_key)
        self._registry.register_handler("browser_upload", self._handle_browser_upload)
        self._registry.register_handler("browser_console", self._handle_browser_console)
        # File view
        self._registry.register_handler("file_view", self._handle_file_view)
        # Search
        self._registry.register_handler("web_search", self._handle_web_search)
        # Deploy
        self._registry.register_handler("deploy_to_vps", self._handle_deploy_to_vps)
        # Schedule
        self._registry.register_handler("schedule_task", self._handle_schedule_task)
        # Image generation
        self._registry.register_handler("image_generate", self._handle_image_generate)
        # Golden paths & delivery
        self._registry.register_handler("get_template", self._handle_get_template)
        self._registry.register_handler("pexels_search", self._handle_pexels_search)
        self._registry.register_handler("create_archive", self._handle_create_archive)
        # Design judge
        self._registry.register_handler("design_judge", self._handle_design_judge)
        # Design RAG search
        self._registry.register_handler("search_design_inspiration", self._handle_search_design_inspiration)
        # Scratchpad
        self._registry.register_handler("update_scratchpad", self._handle_update_scratchpad)


        # ── Document tools ──
        self._registry.register_handler("md_to_pdf", self._handle_md_to_pdf)
        self._registry.register_handler("create_excel", self._handle_create_excel)
        self._registry.register_handler("read_document", self._handle_read_document)
        self._registry.register_handler("create_presentation", self._handle_create_presentation)
        # ── Media tools ──
        self._registry.register_handler("image_edit", self._handle_image_edit)
        self._registry.register_handler("render_diagram", self._handle_render_diagram)
        self._registry.register_handler("generate_chart", self._handle_generate_chart)
        # ── Advanced shell tools ──
        self._registry.register_handler("shell_send", self._handle_shell_send)
        self._registry.register_handler("shell_wait", self._handle_shell_wait)
        self._registry.register_handler("shell_kill", self._handle_shell_kill)
        # ── Advanced browser tools ──
        self._registry.register_handler("browser_fill_form", self._handle_browser_fill_form)
        self._registry.register_handler("browser_move_mouse", self._handle_browser_move_mouse)
        self._registry.register_handler("browser_close", self._handle_browser_close)
        # ── Advanced tools ──
        self._registry.register_handler("expose_port", self._handle_expose_port)
        self._registry.register_handler("speech_to_text", self._handle_speech_to_text)
        self._registry.register_handler("text_to_speech", self._handle_text_to_speech)
        self._registry.register_handler("parallel_map", self._handle_parallel_map)
        # ── Skills system ──
        self._registry.register_handler("read_skill", self._handle_read_skill)
        self._registry.register_handler("list_skills", self._handle_list_skills)
        # ── WebDev scaffolding ──
        self._registry.register_handler("init_project", self._handle_init_project)
        # ── Enhanced search ──
        self._registry.register_handler("search", self._handle_search)

    def get_tools_schema(self) -> list[dict]:
        """Get tools schema for LLM function calling."""
        return self._registry.get_tools_schema()

    async def execute(
        self,
        tool_name: str,
        arguments: Any,
        project_id: str = "",
        user_id: str = "",
    ) -> str:
        """Execute a tool call and return the result as a string."""
        handler = self._registry.get_handler(tool_name)
        if not handler:
            return f"Error: Unknown tool '{tool_name}'. Available tools: {', '.join(self._registry.get_tool_names())}"

        # Parse arguments if string
        if isinstance(arguments, str):
            try:
                arguments = json.loads(arguments)
            except json.JSONDecodeError:
                arguments = {"raw": arguments}

        log_with_data(
            logger, "INFO",
            f"Executing tool: {tool_name}",
            tool=tool_name,
            project_id=project_id,
        )

        try:
            result = await handler(arguments, project_id=project_id, user_id=user_id)
            return str(result) if result is not None else "OK"
        except Exception as e:
            error_report = analyze_error(str(e))
            return f"Tool execution error: {e}\nCategory: {error_report.category.value}\nFixes: {', '.join(error_report.suggested_fixes)}"

    def _get_session(self, session_id: str = "default") -> SandboxSession:
        """Get or create a sandbox session."""
        if session_id not in self._sessions:
            self._sessions[session_id] = SandboxSession(
                session_id=session_id,
                working_dir=self._project_dir,
            )
        return self._sessions[session_id]

    # ── Shell Handlers ────────────────────────────────────────────────────────

    async def _handle_shell_exec(self, args: dict, **kwargs) -> str:
        command = args.get("command", "")
        timeout = args.get("timeout", 30)
        working_dir = args.get("working_dir")
        session = self._get_session(args.get("session", "default"))

        result = await session.execute(command, timeout=timeout, working_dir=working_dir)

        output_parts = []
        if result["stdout"]:
            output_parts.append(result["stdout"])
        if result["stderr"]:
            output_parts.append(f"STDERR: {result['stderr']}")
        output_parts.append(f"[exit code: {result['exit_code']}, {result['elapsed_seconds']}s]")

        return "\n".join(output_parts)

    async def _handle_shell_view(self, args: dict, **kwargs) -> str:
        session = self._get_session(args.get("session", "default"))
        if session._history:
            last = session._history[-1]
            return f"Last command: {last['command']}\nOutput: {last['stdout'][:2000]}"
        return "No commands executed in this session yet."

    async def _handle_ssh_exec(self, args: dict, **kwargs) -> str:
        """Execute a command on a remote server via SSH using paramiko."""
        import paramiko, socket
        host = args.get("host", "")
        command = args.get("command", "")
        username = args.get("username", "root")
        password = args.get("password")
        port = int(args.get("port", 22))
        timeout = int(args.get("timeout", 60))
        key_path = args.get("key_path")
        if not host:
            return "Error: host is required"
        if not command:
            return "Error: command is required"
        try:
            client = paramiko.SSHClient()
            client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
            connect_kwargs = dict(hostname=host, port=port, username=username, timeout=15)
            if key_path:
                connect_kwargs["key_filename"] = key_path
            elif password:
                connect_kwargs["password"] = password
                connect_kwargs["look_for_keys"] = False
                connect_kwargs["allow_agent"] = False
            client.connect(**connect_kwargs)
            stdin, stdout, stderr = client.exec_command(command, timeout=timeout, get_pty=True)
            out = stdout.read().decode("utf-8", errors="replace")
            err = stderr.read().decode("utf-8", errors="replace")
            exit_code = stdout.channel.recv_exit_status()
            client.close()
            result_parts = []
            if out:
                result_parts.append(out)
            if err and err.strip():
                result_parts.append(f"STDERR: {err}")
            result_parts.append(f"[exit code: {exit_code}]")
            return "\n".join(result_parts)
        except paramiko.AuthenticationException:
            return f"Error: SSH authentication failed for {username}@{host}:{port}"
        except (paramiko.SSHException, socket.error) as e:
            return f"Error: SSH connection failed to {host}:{port} — {e}"
        except Exception as e:
            return f"Error: {type(e).__name__}: {e}"

    # ── File Handlers ─────────────────────────────────────────────────────────

    async def _handle_file_read(self, args: dict, **kwargs) -> str:
        path = args.get("path", "")
        is_valid, result = _validate_file_path(path)
        if not is_valid:
            return f"Error: {result}"
        path = result  # use resolved realpath
        start_line = args.get("start_line")
        end_line = args.get("end_line")

        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()

            if start_line or end_line:
                start = (start_line or 1) - 1
                end = end_line if end_line and end_line > 0 else len(lines)
                lines = lines[start:end]

            content = "".join(lines)
            if len(content) > 10000:
                content = content[:10000] + f"\n\n... [truncated, total {len(content)} chars]"
            return content

        except FileNotFoundError:
            return f"Error: File not found: {path}"
        except Exception as e:
            return f"Error reading file: {e}"

    async def _handle_file_write(self, args: dict, **kwargs) -> str:
        path = args.get("path", "")
        is_valid, result = _validate_file_path(path)
        if not is_valid:
            return f"Error: {result}"
        path = result
        content = args.get("content", "")

        try:
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            return f"File written: {path} ({len(content)} chars)"
        except Exception as e:
            return f"Error writing file: {e}"

    async def _handle_file_edit(self, args: dict, **kwargs) -> str:
        path = args.get("path", "")
        is_valid, result = _validate_file_path(path)
        if not is_valid:
            return f"Error: {result}"
        path = result
        edits = args.get("edits", [])

        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()

            original = content
            changes = 0
            for edit in edits:
                find = edit.get("find", "")
                replace = edit.get("replace", "")
                replace_all = edit.get("all", False)

                if find not in content:
                    return f"Error: Text not found in file: '{find[:50]}...'"

                if replace_all:
                    count = content.count(find)
                    content = content.replace(find, replace)
                    changes += count
                else:
                    content = content.replace(find, replace, 1)
                    changes += 1

            with open(path, "w", encoding="utf-8") as f:
                f.write(content)

            return f"File edited: {path} ({changes} changes applied)"

        except FileNotFoundError:
            return f"Error: File not found: {path}"
        except Exception as e:
            return f"Error editing file: {e}"

    async def _handle_file_append(self, args: dict, **kwargs) -> str:
        path = args.get("path", "")
        is_valid, result = _validate_file_path(path)
        if not is_valid:
            return f"Error: {result}"
        path = result
        content = args.get("content", "")

        try:
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "a", encoding="utf-8") as f:
                f.write(content)
            return f"Content appended to: {path}"
        except Exception as e:
            return f"Error appending to file: {e}"

    # ── Match Handlers ────────────────────────────────────────────────────────

    async def _handle_glob(self, args: dict, **kwargs) -> str:
        import glob
        pattern = args.get("pattern", "")
        # Security: validate that the glob base path is within allowed dirs
        base_path = pattern.split("*")[0].rstrip("/") if "*" in pattern else os.path.dirname(pattern)
        if base_path:
            valid, msg = _validate_file_path(base_path)
            if not valid:
                return f"Error: {msg}"
        matches = sorted(glob.glob(pattern, recursive=True))
        if not matches:
            return f"No files matching: {pattern}"
        # FIX 3: Filter glob results — validate each path is within allowed dirs
        safe_matches = []
        for m in matches[:100]:
            valid, _ = _validate_file_path(m)
            if valid:
                safe_matches.append(m)
        return "\n".join(safe_matches) if safe_matches else "No accessible files matching pattern"

    async def _handle_grep(self, args: dict, **kwargs) -> str:
        import re as re_module
        pattern = args.get("pattern", "")
        scope = args.get("scope", "")
        context_lines = args.get("context_lines", 0)

        # Security: validate scope base path
        scope_base = scope.split("*")[0].rstrip("/") if "*" in scope else os.path.dirname(scope)
        if scope_base:
            valid, msg = _validate_file_path(scope_base)
            if not valid:
                return f"Error: {msg}"

        import glob
        files = glob.glob(scope, recursive=True)
        results = []

        for filepath in files[:50]:  # Limit files
            # FIX 3: Validate each file path before opening
            valid, _ = _validate_file_path(filepath)
            if not valid:
                continue
            try:
                with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                    lines = f.readlines()
                for i, line in enumerate(lines):
                    if re_module.search(pattern, line):
                        start = max(0, i - context_lines)
                        end = min(len(lines), i + context_lines + 1)
                        snippet = "".join(lines[start:end])
                        results.append(f"{filepath}:{i+1}\n{snippet}")
            except Exception:
                continue

        if not results:
            return f"No matches for '{pattern}' in {scope}"
        return "\n---\n".join(results[:20])

    # ── Communication Handlers ────────────────────────────────────────────────

    async def _handle_message(self, args: dict, **kwargs) -> str:
        msg_type = args.get("type", "info")
        text = args.get("text", "")
        attachments = args.get("attachments", [])
        # This will be intercepted by the agent loop for WebSocket delivery
        return json.dumps({
            "type": msg_type,
            "text": text,
            "attachments": attachments,
        }, ensure_ascii=False)

    async def _handle_plan(self, args: dict, **kwargs) -> str:
        action = args.get("action", "update")
        # This will be intercepted by the agent loop for plan management
        return json.dumps({
            "action": action,
            "goal": args.get("goal"),
            "phases": args.get("phases"),
            "current_phase_id": args.get("current_phase_id"),
            "next_phase_id": args.get("next_phase_id"),
        })

    # ── Browser Handlers ─────────────────────────────────────────────────────

    async def _get_browser(self):
        """Lazy-initialize and return the BrowserWorker. Auto-restart on crash."""
        if self._browser is None:
            from workers.browser.worker import BrowserWorker
            self._browser = BrowserWorker(
                headless=True,
                screenshots_dir=os.path.join(self._project_dir, "screenshots"),
            )
            return self._browser
        # Check if browser is still alive; reinitialize if crashed
        try:
            if not self._browser._initialized or self._browser._browser is None:
                raise RuntimeError("Browser not initialized")
            # Quick health check — if page is closed, reinit
            if self._browser._page is not None and self._browser._page.is_closed():
                raise RuntimeError("Page is closed")
        except Exception:
            logger.warning("Browser crashed or disconnected — reinitializing...")
            try:
                await self._browser.close()
            except Exception:
                pass
            self._browser._initialized = False
            self._browser._browser = None
            self._browser._context = None
            self._browser._page = None
            self._browser._playwright = None
        return self._browser

    async def _handle_browser_navigate(self, args: dict, **kwargs) -> str:
        url = args.get("url", "")
        if not url:
            return "Error: 'url' parameter is required for browser_navigate"

        # Bug #10 fix: Chromium blocks file:/// in headless mode.
        # Spin up a temporary HTTP server and rewrite the URL.
        if url.startswith("file:///"):
            import urllib.parse
            file_path = urllib.parse.unquote(url[7:])  # strip file:///
            serve_dir = os.path.dirname(file_path) or "/root/workspace"
            filename = os.path.basename(file_path) or "index.html"

            # Start a one-off HTTP server if not already running for this dir
            if not hasattr(self, '_file_server_proc') or self._file_server_proc is None or self._file_server_proc.poll() is not None:
                import subprocess as _sp
                self._file_server_proc = _sp.Popen(
                    ["python3", "-m", "http.server", "8079", "--directory", serve_dir],
                    stdout=_sp.DEVNULL, stderr=_sp.DEVNULL,
                )
                await asyncio.sleep(0.5)  # let server start
                logger.info(f"Started file server on :8079 for {serve_dir}")

            url = f"http://localhost:8079/{filename}"
            logger.info(f"Rewrote file:/// URL to {url}")

        browser = await self._get_browser()
        result = await browser.navigate(url)
        if "error" in result:
            return f"Browser navigation error: {result['error']}"
        # Return structured info (screenshot path, elements, content)
        elements_text = "\n".join(
            f"  [{e['index']}] <{e['tag']}>{e.get('text', '')[:40]}</{e['tag']}>"
            for e in result.get("elements", [])[:30]
        )
        b64 = result.get("screenshot_b64", "")
        b64_line = f"\nscreenshot_b64:{b64}" if b64 else ""
        return (
            f"Navigated to: {result.get('url', url)}\n"
            f"Title: {result.get('title', '')}\n"
            f"Status: {result.get('status', '')}\n"
            f"Screenshot: {result.get('screenshot', '')}\n"
            f"Interactive elements ({result.get('element_count', 0)}):\n{elements_text}\n\n"
            f"Page content (first 3000 chars):\n{result.get('content', '')[:3000]}"
            f"{b64_line}"
        )

    async def _handle_browser_view(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        result = await browser.get_view()
        elements_text = "\n".join(
            f"  [{e['index']}] <{e['tag']}>{e.get('text', '')[:40]}</{e['tag']}>"
            for e in result.get("elements", [])[:30]
        )
        b64 = result.get("screenshot_b64", "")
        b64_line = f"\nscreenshot_b64:{b64}" if b64 else ""
        return (
            f"URL: {result.get('url', '')}\n"
            f"Title: {result.get('title', '')}\n"
            f"Screenshot: {result.get('screenshot', '')}\n"
            f"Elements:\n{elements_text}\n\n"
            f"Content:\n{result.get('content', '')[:3000]}"
            f"{b64_line}"
        )

    async def _handle_browser_click(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        index = args.get("index")
        x = args.get("coordinate_x")
        y = args.get("coordinate_y")
        result = await browser.click(index=index, x=x, y=y)
        if "error" in result:
            return f"Click error: {result['error']}"
        b64 = result.get("screenshot_b64", "")
        b64_line = f"\nscreenshot_b64:{b64}" if b64 else ""
        return f"Clicked successfully. URL: {result.get('url', '')}\nScreenshot: {result.get('screenshot', '')}{b64_line}"

    async def _handle_browser_input(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        text = args.get("text", "")
        index = args.get("index")
        x = args.get("coordinate_x")
        y = args.get("coordinate_y")
        press_enter = args.get("press_enter", False)
        result = await browser.input_text(text, index=index, x=x, y=y, press_enter=press_enter)
        if "error" in result:
            return f"Input error: {result['error']}"
        b64 = result.get("screenshot_b64", "")
        b64_line = f"\nscreenshot_b64:{b64}" if b64 else ""
        return f"Text entered successfully. Screenshot: {result.get('screenshot', '')}{b64_line}"

    async def _handle_browser_scroll(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        direction = args.get("direction", "down")
        to_end = args.get("to_end", False)
        result = await browser.scroll(direction=direction, to_end=to_end)
        if "error" in result:
            return f"Scroll error: {result['error']}"
        b64 = result.get("screenshot_b64", "")
        b64_line = f"\nscreenshot_b64:{b64}" if b64 else ""
        return f"Scrolled {direction}. Screenshot: {result.get('screenshot', '')}{b64_line}"

    async def _handle_browser_select(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        index = args.get("index", 0)
        option_index = args.get("option_index", 0)
        # Map index to element selector from cache
        if index < len(browser._elements_cache):
            el = browser._elements_cache[index]
            selector = el.get("selector", f"select:nth-of-type({index + 1})")
            result = await browser.select_option(selector, str(option_index))
        else:
            result = {"error": f"Element index {index} not found in cache"}
        if "error" in result:
            return f"Select error: {result['error']}"
        return f"Option selected. Screenshot: {result.get('screenshot', '')}"

    async def _handle_browser_find(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        keyword = args.get("keyword", "")
        result = await browser.find_text(keyword)
        if "error" in result:
            return f"Find error: {result['error']}"
        if result.get("found"):
            matches = "\n".join(f"  ...{m}..." for m in result.get("matches", []))
            return f"Found {result['count']} matches for '{keyword}':\n{matches}"
        return f"No matches found for '{keyword}'"

    async def _handle_browser_save_image(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        x = args.get("coordinate_x", 0)
        y = args.get("coordinate_y", 0)
        save_dir = args.get("save_dir", self._project_dir)
        base_name = args.get("base_name", "image")
        result = await browser.save_image(x, y, save_dir, base_name)
        if "error" in result:
            return f"Save image error: {result['error']}"
        return f"Image saved: {result.get('path', '')}"

    async def _handle_browser_press_key(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        key = args.get("key", "")
        await browser.initialize()
        try:
            await browser._page.keyboard.press(key)
            await asyncio.sleep(0.3)
            screenshot = await browser._take_screenshot()
            return f"Key '{key}' pressed. Screenshot: {screenshot}"
        except Exception as e:
            return f"Key press error: {e}"

    async def _handle_browser_upload(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        index = args.get("index", 0)
        path = args.get("path", "")
        # Security: validate file path before uploading
        valid, resolved = _validate_file_path(path)
        if not valid:
            return f"Error: {resolved}"
        path = resolved
        await browser.initialize()
        try:
            if index < len(browser._elements_cache):
                el = browser._elements_cache[index]
                selector = el.get("selector", f"input[type=file]:nth-of-type({index + 1})")
                file_input = await browser._page.query_selector(selector)
                if file_input:
                    await file_input.set_input_files(path)
                    return f"File uploaded: {path}"
            return f"Error: File input element at index {index} not found"
        except Exception as e:
            return f"Upload error: {e}"

    async def _handle_browser_console(self, args: dict, **kwargs) -> str:
        browser = await self._get_browser()
        js = args.get("javascript", "")
        result = await browser.execute_js(js)
        if "error" in result:
            return f"JS error: {result['error']}"
        return f"JS result: {result.get('result', 'undefined')}"

    # ── File View Handler ────────────────────────────────────────────────────

    async def _handle_file_view(self, args: dict, **kwargs) -> str:
        path = args.get("path", "")
        # Security: validate path is within allowed directories
        valid, resolved = _validate_file_path(path)
        if not valid:
            return f"Error: {resolved}"
        path = resolved
        if not os.path.exists(path):
            return f"Error: File not found: {path}"
        ext = os.path.splitext(path)[1].lower()
        if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"):
            return f"Image file: {path} ({os.path.getsize(path)} bytes)"
        elif ext == ".pdf":
            try:
                import subprocess
                result = subprocess.run(
                    ["pdftotext", path, "-"],
                    capture_output=True, text=True, timeout=10
                )
                return result.stdout[:5000] if result.stdout else f"PDF file: {path}"
            except Exception:
                return f"PDF file: {path} ({os.path.getsize(path)} bytes)"
        else:
            return await self._handle_file_read(args, **kwargs)

    # ── Web Search Handler ───────────────────────────────────────────────────

    async def _handle_web_search(self, args: dict, **kwargs) -> str:
        queries = args.get("queries", [])
        search_type = args.get("type", "info")
        if not queries:
            return "Error: 'queries' parameter is required"

        # FAST-SKIP: For well-known technologies, skip search and tell agent to use knowledge
        KNOWN_TECH_KEYWORDS = [
            "docker-compose", "dockerfile", "docker compose", "nginx", "apache",
            "systemd", "cron", "crontab", "makefile", "bash script", "shell script",
            "python script", "pip install", "requirements.txt", "package.json",
            "tsconfig", "eslint", "prettier", "webpack", "vite",
            "terraform", "kubernetes", "k8s", "helm", "github actions",
            ".env", "environment variables", "gitignore", ".gitignore",
            "sql query", "create table", "postgresql config", "redis config",
            "ssh config", "firewall", "iptables", "ufw",
        ]
        query_lower = " ".join(queries).lower()
        for kw in KNOWN_TECH_KEYWORDS:
            if kw in query_lower:
                logger.info(f"web_search FAST-SKIP: '{kw}' is a known technology, skipping search")
                return f"Search skipped — {kw} is a well-known technology. Use your built-in knowledge to create the file directly. You know the correct syntax and best practices."

        import os
        has_api_keys = any([
            os.getenv("TAVILY_API_KEY", ""),
            os.getenv("SERPER_API_KEY", ""),
            os.getenv("EXA_API_KEY", ""),
            os.getenv("BRAVE_API_KEY", ""),
        ])

        # Try SearchWorker only if API keys are configured
        if has_api_keys:
            try:
                from workers.search.worker import SearchWorker
                worker = SearchWorker()
                result = await worker.search(
                    queries=queries,
                    search_type=search_type,
                    max_results=10,
                )
                formatted = []
                for r in result.get("results", []):
                    title = r.get("title", "")
                    url = r.get("url", "")
                    snippet = r.get("snippet", r.get("content", ""))[:500]
                    formatted.append(f"**{title}**\n{url}\n{snippet}")
                if formatted:
                    return "\n\n---\n\n".join(formatted)
            except Exception as e:
                logger.warning(f"SearchWorker failed: {e}")

        # Fast DuckDuckGo fallback using JSON API
        import aiohttp
        results = []
        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as http:
            for query in queries[:2]:  # Max 2 queries for speed
                try:
                    # DuckDuckGo instant answer API (fast, no scraping)
                    url = f"https://api.duckduckgo.com/?q={query}&format=json&no_html=1&skip_disambig=1"
                    async with http.get(url) as resp:
                        if resp.status == 200:
                            data = await resp.json(content_type=None)
                            abstract = data.get("AbstractText", "")
                            abstract_url = data.get("AbstractURL", "")
                            answer = data.get("Answer", "")
                            # Collect related topics
                            topics = []
                            for t in data.get("RelatedTopics", [])[:5]:
                                if isinstance(t, dict) and "Text" in t:
                                    topics.append(f"- {t['Text'][:200]}")
                            
                            parts = []
                            if answer:
                                parts.append(f"Answer: {answer}")
                            if abstract:
                                parts.append(f"{abstract}\nSource: {abstract_url}")
                            if topics:
                                parts.append("Related:\n" + "\n".join(topics))
                            if parts:
                                results.append(f"Query: {query}\n" + "\n".join(parts))
                            else:
                                results.append(f"Query: {query}\nNo instant answer available. Use your knowledge to answer.")
                except Exception as e:
                    results.append(f"Query: {query} — Search unavailable: {e}")

        if not results:
            return "Search APIs unavailable. Use your built-in knowledge to answer the question. If you need specific current data (prices, news), tell the user that real-time search is temporarily unavailable."
        return "\n\n---\n\n".join(results)
    # ── Deploy Handlers ──────────────────────────────────────────────────────

    async def _handle_deploy_to_vps(self, args: dict, **kwargs) -> str:
        source_dir = args.get("source_dir", "")
        domain = args.get("domain", "")
        server_host = args.get("server_host", "")
        server_user = args.get("server_user", "root")
        deploy_type = args.get("deploy_type", "static")

        if not source_dir or not domain:
            return "Error: source_dir and domain are required"

        # This is a placeholder — actual deploy logic uses SSH worker
        return json.dumps({
            "type": "deploy_to_vps",
            "source_dir": source_dir,
            "domain": domain,
            "server_host": server_host,
            "server_user": server_user,
            "deploy_type": deploy_type,
        })

    async def _handle_schedule_task(self, args: dict, **kwargs) -> str:
        return json.dumps({"type": "schedule_task", **args})

    # ── Image Generation Handler ────────────────────────────────────────────

    async def _handle_image_generate(self, args: dict, **kwargs) -> str:
        """Generate images using DALL-E 3 or FLUX."""
        from workers.image_gen import get_image_generator

        prompt = args.get("prompt", "")
        if not prompt:
            return "Error: prompt is required"

        style = args.get("style", "photorealistic")
        size = args.get("size", "1024x1024")
        quality = args.get("quality", "standard")
        n = min(args.get("n", 1), 4)
        project_id = kwargs.get("project_id", "")
        save_dir = args.get("save_dir") or os.path.join(self._project_dir, "images")

        generator = get_image_generator()
        result = await generator.generate(
            prompt=prompt,
            style=style,
            size=size,
            quality=quality,
            n=n,
            project_id=project_id,
            save_dir=save_dir,
        )

        if result["success"]:
            paths = [img["path"] for img in result["images"]]
            return json.dumps({
                "success": True,
                "images": paths,
                "provider": result["provider"],
                "cost": result["cost"],
                "elapsed_seconds": result["elapsed_seconds"],
            })
        else:
            return f"Image generation failed: {result.get('error', 'Unknown error')}"

    # ── Golden Paths & Archive Handlers ──────────────────────────────────────

    async def _handle_get_template(self, args: dict, **kwargs) -> str:
        """Get a golden path template for landing page generation.
        
        Returns one of the 7 premium blueprint HTML scaffolds based on the requested template_type.
        """
        from core.golden_paths import get_template, list_templates

        template_type = args.get("template_type", "").lower().strip()

        if template_type == "list":
            templates = list_templates()
            return json.dumps({"templates": templates}, ensure_ascii=False)

        # Normalize aliases to known blueprint types
        alias_map = {
            "barber": "dark_luxury", "барбершоп": "dark_luxury", "nightclub": "dark_luxury", "jewelry": "dark_luxury", "auto": "dark_luxury",
            "ресторан": "warm_editorial", "кафе": "warm_editorial", "еда": "warm_editorial", "restaurant": "warm_editorial", "bakery": "warm_editorial",
            "tech": "clean_tech", "it": "clean_tech", "стартап": "clean_tech", "saas": "clean_tech", "b2b": "clean_tech",
            "фитнес": "bold_energy", "спорт": "bold_energy", "gym": "bold_energy", "fitness": "bold_energy", "event": "bold_energy",
            "стоматология": "soft_wellness", "dental": "soft_wellness", "врач": "soft_wellness", "spa": "soft_wellness", "beauty": "soft_wellness", "medical": "soft_wellness",
            "фото": "japandi_minimal", "portfolio": "japandi_minimal", "architecture": "japandi_minimal", "interior": "japandi_minimal",
            "юрист": "neobrutalist", "agency": "neobrutalist", "web3": "neobrutalist", "creative": "neobrutalist",
        }
        
        resolved_type = template_type
        if template_type in alias_map:
            resolved_type = alias_map[template_type]
            
        # Default to clean_tech if unknown
        valid_blueprints = ["dark_luxury", "warm_editorial", "clean_tech", "bold_energy", "soft_wellness", "japandi_minimal", "neobrutalist"]
        if resolved_type not in valid_blueprints:
            resolved_type = "clean_tech"

        scaffold_path = f"/root/arcane/templates/blueprints/{resolved_type}.html"
        scaffold_html = ""
        try:
            with open(scaffold_path, "r", encoding="utf-8") as f:
                scaffold_html = f.read()
            logger.info(f"Loaded blueprint scaffold ({len(scaffold_html)} chars) for template_type='{resolved_type}'")
        except Exception as e:
            logger.error(f"Failed to load blueprint scaffold: {e}")

        result = {
            "scaffold_html": scaffold_html if scaffold_html else "ERROR: scaffold not found at " + scaffold_path,
            "template_type": resolved_type,
            "instructions": (
                "This is your MANDATORY blueprint scaffold. Do NOT generate from scratch. "
                "Follow the <landing_quality_standard> in your system prompt: "
                "1. Replace {{PLACEHOLDERS}} with real client content. "
                "2. Inject REAL photos (from Pexels API) into the image src attributes. "
                "3. DO NOT add generic AI effects like particles or glassmorphism. "
                "4. Keep the blueprint's layout, typography, and structure intact."
            ),
        }

        return json.dumps(result, ensure_ascii=False)

    async def _handle_create_archive(self, args: dict, **kwargs) -> str:
        """Package project files into a downloadable ZIP archive."""
        from core.golden_paths import create_delivery_archive

        project_dir = args.get("project_dir", "")
        if not project_dir:
            project_dir = self._project_dir

        project_name = args.get("project_name", "arcane-project")
        exclude = args.get("exclude", None)

        try:
            archive_path = create_delivery_archive(
                project_dir=project_dir,
                project_name=project_name,
                include_readme=args.get("include_readme", True),
                include_deploy_instructions=args.get("include_deploy_guide", True),
                exclude_patterns=exclude,
            )
            size_mb = os.path.getsize(archive_path) / (1024 * 1024)
            return json.dumps({
                "success": True,
                "archive_path": archive_path,
                "size_mb": round(size_mb, 2),
                "project_name": project_name,
            })
        except Exception as e:
            return f"Error creating archive: {e}"

    async def _handle_pexels_search(self, args: dict, **kwargs) -> str:
        """Search for high-quality stock photos. Uses Pexels API if key available, falls back to Unsplash source URLs."""
        import json
        import os
        import urllib.parse
        
        query = args.get("query", "")
        per_page = args.get("per_page", 5)
        
        if not query:
            return "Error: query is required"
        
        api_key = os.environ.get("PEXELS_API_KEY", "")
        
        # Try Pexels API first if key is available
        if api_key:
            try:
                import aiohttp
                async with aiohttp.ClientSession() as session:
                    headers = {"Authorization": api_key}
                    encoded_query = urllib.parse.quote(query)
                    url = f"https://api.pexels.com/v1/search?query={encoded_query}&per_page={per_page}&orientation=landscape"
                    async with session.get(url, headers=headers, timeout=aiohttp.ClientTimeout(total=10)) as response:
                        if response.status == 200:
                            data = await response.json()
                            photos = []
                            for photo in data.get("photos", []):
                                photos.append({
                                    "id": photo.get("id"),
                                    "src": photo.get("src", {}).get("large2x", photo.get("src", {}).get("original", "")),
                                    "src_medium": photo.get("src", {}).get("medium", ""),
                                    "src_landscape": photo.get("src", {}).get("landscape", ""),
                                    "alt": photo.get("alt", query),
                                    "photographer": photo.get("photographer", ""),
                                    "avg_color": photo.get("avg_color", "")
                                })
                            if photos:
                                return json.dumps({"source": "pexels", "photos": photos}, ensure_ascii=False)
            except Exception as e:
                logger.warning(f"Pexels API failed: {e}, falling back to Unsplash")
        
        # Fallback: curated stock photo collections by category
        # These are reliable direct Unsplash URLs that don't require API keys
        curated_photos = {
            # Business / Office
            "business": [
                {"src": "https://images.unsplash.com/photo-1497366216548-37526070297c?w=1200&h=800&fit=crop", "alt": "Modern office"},
                {"src": "https://images.unsplash.com/photo-1497215842964-222b430dc094?w=1200&h=800&fit=crop", "alt": "Office workspace"},
                {"src": "https://images.unsplash.com/photo-1556761175-4b46a572b786?w=1200&h=800&fit=crop", "alt": "Team meeting"},
            ],
            # Food / Restaurant
            "food": [
                {"src": "https://images.unsplash.com/photo-1414235077428-338989a2e8c0?w=1200&h=800&fit=crop", "alt": "Fine dining"},
                {"src": "https://images.unsplash.com/photo-1517248135467-4c7edcad34c4?w=1200&h=800&fit=crop", "alt": "Restaurant interior"},
                {"src": "https://images.unsplash.com/photo-1544025162-d76694265947?w=1200&h=800&fit=crop", "alt": "Gourmet dish"},
            ],
            # Barbershop / Salon
            "barber": [
                {"src": "https://images.unsplash.com/photo-1503951914875-452162b0f3f1?w=1200&h=800&fit=crop", "alt": "Barbershop"},
                {"src": "https://images.unsplash.com/photo-1599351431202-1e0f0137899a?w=1200&h=800&fit=crop", "alt": "Haircut"},
                {"src": "https://images.unsplash.com/photo-1621605815971-fbc98d665033?w=1200&h=800&fit=crop", "alt": "Barber tools"},
            ],
            # Fitness / Sports
            "fitness": [
                {"src": "https://images.unsplash.com/photo-1534438327276-14e5300c3a48?w=1200&h=800&fit=crop", "alt": "Gym"},
                {"src": "https://images.unsplash.com/photo-1571019613454-1cb2f99b2d8b?w=1200&h=800&fit=crop", "alt": "Workout"},
                {"src": "https://images.unsplash.com/photo-1517836357463-d25dfeac3438?w=1200&h=800&fit=crop", "alt": "Training"},
            ],
            # Medical / Wellness
            "medical": [
                {"src": "https://images.unsplash.com/photo-1519494026892-80bbd2d6fd0d?w=1200&h=800&fit=crop", "alt": "Medical clinic"},
                {"src": "https://images.unsplash.com/photo-1576091160399-112ba8d25d1d?w=1200&h=800&fit=crop", "alt": "Doctor"},
                {"src": "https://images.unsplash.com/photo-1629909613654-28e377c37b09?w=1200&h=800&fit=crop", "alt": "Dental office"},
            ],
            # Technology / SaaS
            "tech": [
                {"src": "https://images.unsplash.com/photo-1518770660439-4636190af475?w=1200&h=800&fit=crop", "alt": "Technology"},
                {"src": "https://images.unsplash.com/photo-1551434678-e076c223a692?w=1200&h=800&fit=crop", "alt": "Developer team"},
                {"src": "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=1200&h=800&fit=crop", "alt": "Dashboard"},
            ],
            # Beauty / Spa
            "beauty": [
                {"src": "https://images.unsplash.com/photo-1560750588-73b555dce5d1?w=1200&h=800&fit=crop", "alt": "Spa"},
                {"src": "https://images.unsplash.com/photo-1570172619644-dfd03ed5d881?w=1200&h=800&fit=crop", "alt": "Beauty treatment"},
                {"src": "https://images.unsplash.com/photo-1522337360788-8b13dee7a37e?w=1200&h=800&fit=crop", "alt": "Skincare"},
            ],
            # Architecture / Interior
            "architecture": [
                {"src": "https://images.unsplash.com/photo-1487958449943-2429e8be8625?w=1200&h=800&fit=crop", "alt": "Architecture"},
                {"src": "https://images.unsplash.com/photo-1600607687939-ce8a6c25118c?w=1200&h=800&fit=crop", "alt": "Interior design"},
                {"src": "https://images.unsplash.com/photo-1618221195710-dd6b41faaea6?w=1200&h=800&fit=crop", "alt": "Modern interior"},
            ],
        }
        
        # Find best matching category
        query_lower = query.lower()
        best_category = "business"  # default
        category_keywords = {
            "food": ["food", "restaurant", "cafe", "кафе", "ресторан", "еда", "dish", "cuisine", "menu", "chef", "cook", "bakery", "wine"],
            "barber": ["barber", "haircut", "beard", "salon", "стрижк", "борода", "барбер", "hair"],
            "fitness": ["fitness", "gym", "sport", "workout", "фитнес", "спорт", "тренировк", "training"],
            "medical": ["medical", "doctor", "dental", "clinic", "врач", "стоматолог", "клиник", "health", "spa", "wellness"],
            "tech": ["tech", "software", "saas", "startup", "code", "developer", "IT", "стартап", "dashboard"],
            "beauty": ["beauty", "spa", "cosmetic", "skincare", "красот", "косметик", "маникюр", "nail"],
            "architecture": ["architect", "interior", "design", "portfolio", "архитектур", "интерьер", "дизайн"],
            "business": ["business", "office", "corporate", "team", "бизнес", "офис"],
        }
        
        for cat, keywords in category_keywords.items():
            if any(kw in query_lower for kw in keywords):
                best_category = cat
                break
        
        selected = curated_photos.get(best_category, curated_photos["business"])
        photos = []
        for i in range(min(per_page, len(selected))):
            photo = selected[i]
            photos.append({
                "src": photo["src"],
                "src_medium": photo["src"].replace("w=1200&h=800", "w=600&h=400"),
                "src_landscape": photo["src"].replace("w=1200&h=800", "w=1200&h=627"),
                "alt": photo["alt"],
                "photographer": "Unsplash",
            })
        
        return json.dumps({
            "source": "curated_unsplash",
            "photos": photos,
            "note": f"Using curated Unsplash photos for category '{best_category}'. For dynamic results, set PEXELS_API_KEY."
        }, ensure_ascii=False)

    # ── Scratchpad Handler ────────────────────────────────────────────────────

    async def _handle_update_scratchpad(self, args: dict, **kwargs) -> str:
        """Save a key-value pair to the agent's scratchpad."""
        key = args.get("key", "")
        value = args.get("value", "")

        if not key:
            return "Error: key is required"

        # Store in module-level fallback; agent_loop will also update its own scratchpad
        if not hasattr(self, "_scratchpad_store"):
            self._scratchpad_store = {}
        self._scratchpad_store[key] = value
        return f"Scratchpad updated: {key} = {value}"

    # ── Design Judge Handler ─────────────────────────────────────────────────

    async def _handle_design_judge(self, args: dict, **kwargs) -> str:
        """Evaluate a generated website using Vision Judge v2 (Playwright screenshot + Vision API)."""
        from workers.design_judge_v2 import get_vision_judge
        from config.settings import get_config
        judge = get_vision_judge(get_config())
        html_path = args.get("html_path", "")
        screenshot_path = args.get("screenshot_path", "")
        context = args.get("context", "")
        model = args.get("model", "google/gemini-2.5-flash")
        include_mobile = args.get("include_mobile", True)
        if html_path:
            # Full pipeline: Playwright screenshot + Vision evaluation
            result = await judge.evaluate_html(
                html_path=html_path,
                context=context,
                model=model,
                include_mobile=include_mobile,
            )
        elif screenshot_path:
            # Direct screenshot evaluation (if already have a screenshot)
            result = await judge.evaluate_screenshot(screenshot_path, context, model)
        else:
            return json.dumps({"success": False, "error": "html_path is required. Provide the path to the generated HTML file."})
        # Format output for the agent with clear instructions
        if result.get("success"):
            output = {
                "success": True,
                "overall_score": result.get("overall_score", 0),
                "tier": result.get("tier", "UNKNOWN"),
                "verdict": result.get("verdict", ""),
                "scores": result.get("scores", {}),
                "strengths": result.get("strengths", []),
                "critical_issues": result.get("critical_issues", []),
                "fix_instructions": result.get("fix_instructions", []),
                "mobile_assessment": result.get("mobile_assessment", ""),
                "evaluation_method": result.get("evaluation_method", ""),
                "model_used": result.get("model_used", ""),
                "elapsed_seconds": result.get("elapsed_seconds", 0),
            }
            # Add agent guidance
            score = result.get("overall_score", 0)
            if score < 7.0:
                output["agent_action"] = "MUST_FIX: Score below 7.0. Apply ALL fix_instructions, then re-evaluate."
            elif score < 8.0:
                output["agent_action"] = "SHOULD_FIX: Score is good but not premium. Apply fix_instructions for critical_issues, then re-evaluate."
            else:
                output["agent_action"] = "APPROVED: Score is premium quality. Minor polish only if time permits."
        else:
            output = result
        return json.dumps(output, ensure_ascii=False)

    # ── Design RAG Search Handler ────────────────────────────────────────────

    async def _handle_search_design_inspiration(self, args: dict, **kwargs) -> str:
        """Search curated design references for inspiration before generating landing pages."""
        import json
        from workers.design_rag import get_design_rag
        from config.settings import get_config

        config = get_config()
        rag = get_design_rag(config)

        query = args.get("query", "")
        if not query:
            return json.dumps({"error": "query is required"})

        result = await rag.search(
            query=query,
            style=args.get("style"),
            mood=args.get("mood"),
            industry=args.get("industry"),
            min_tier=args.get("min_tier", "A"),
            limit=args.get("limit", 8),
            diversity=True,
        )

        # Format output for the agent
        if result.get("references"):
            output = {
                "status": "ok",
                "query": result["query"],
                "suggested_blueprint": result["suggested_blueprint"],
                "total_in_db": result["total_found"],
                "references_returned": result["returned"],
                "instruction": (
                    "Study these reference designs carefully. "
                    "Note their color palettes, typography choices, layout patterns, "
                    "hero treatments, and overall mood. Use the suggested_blueprint as "
                    "your starting template, then adapt it based on these references. "
                    "DO NOT copy — synthesize the best patterns into something original."
                ),
                "references": result["references"],
            }
        else:
            output = {
                "status": "no_results",
                "query": query,
                "suggested_blueprint": result.get("suggested_blueprint", "clean_tech"),
                "instruction": (
                    "No exact matches found. Use the suggested_blueprint template "
                    "and apply premium design principles: generous whitespace, "
                    "intentional typography hierarchy, real photography, and "
                    "sophisticated color palette."
                ),
            }

        return json.dumps(output, ensure_ascii=False)

    # ═══════════════════════════════════════════════════════════════════════════
    # NEW HANDLERS — 17 tools added in Manus-style upgrade
    # ═══════════════════════════════════════════════════════════════════════════

    # ── Document Tools ───────────────────────────────────────────────────────

    async def _handle_md_to_pdf(self, args: dict, **kwargs) -> str:
        """Convert Markdown file to PDF."""
        input_path = args.get("input_file", "")
        output_path = args.get("output_file", "")
        if not input_path:
            return "Error: input_file is required"
        valid, resolved = _validate_file_path(input_path)
        if not valid:
            return f"Error: {resolved}"
        input_path = resolved
        if not os.path.exists(input_path):
            return f"Error: File not found: {input_path}"
        if not output_path:
            output_path = input_path.rsplit(".", 1)[0] + ".pdf"
        valid2, resolved2 = _validate_file_path(output_path)
        if not valid2:
            return f"Error: {resolved2}"
        output_path = resolved2
        try:
            result = await execute_sandboxed(
                command=f"weasyprint '{input_path}' '{output_path}'",
                working_dir=os.path.dirname(input_path),
                timeout=60,
            )
            if os.path.exists(output_path):
                size_kb = os.path.getsize(output_path) / 1024
                return json.dumps({"success": True, "path": output_path, "size_kb": round(size_kb, 1)})
            return f"Error: PDF generation failed. {result.get('stderr', '')}"
        except Exception as e:
            return f"Error: {e}"

    async def _handle_create_excel(self, args: dict, **kwargs) -> str:
        """Create an Excel file from structured data."""
        import openpyxl
        output_path = args.get("output_file", os.path.join(self._project_dir, "output.xlsx"))
        valid, resolved = _validate_file_path(output_path)
        if not valid:
            return f"Error: {resolved}"
        output_path = resolved
        sheets = args.get("sheets", [])
        if not sheets:
            headers = args.get("headers", [])
            rows = args.get("rows", [])
            sheets = [{"name": "Sheet1", "headers": headers, "rows": rows}]
        try:
            wb = openpyxl.Workbook()
            for i, sheet_data in enumerate(sheets):
                ws = wb.active if i == 0 else wb.create_sheet()
                ws.title = sheet_data.get("name", f"Sheet{i+1}")
                headers = sheet_data.get("headers", [])
                if headers:
                    ws.append(headers)
                for row in sheet_data.get("rows", []):
                    ws.append(row)
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            wb.save(output_path)
            size_kb = os.path.getsize(output_path) / 1024
            return json.dumps({"success": True, "path": output_path, "size_kb": round(size_kb, 1)})
        except Exception as e:
            return f"Error creating Excel: {e}"

    async def _handle_read_document(self, args: dict, **kwargs) -> str:
        """Read content from PDF, DOCX, XLSX, or text files."""
        path = args.get("path", "")
        if not path:
            return "Error: path is required"
        valid, resolved = _validate_file_path(path)
        if not valid:
            return f"Error: {resolved}"
        path = resolved
        if not os.path.exists(path):
            return f"Error: File not found: {path}"
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == ".pdf":
                import pdfplumber
                text_parts = []
                with pdfplumber.open(path) as pdf:
                    for i, page in enumerate(pdf.pages[:50]):
                        text_parts.append(f"--- Page {i+1} ---\n{page.extract_text() or ''}")
                return "\n".join(text_parts)[:10000]
            elif ext in (".docx", ".doc"):
                import docx
                doc = docx.Document(path)
                return "\n".join(p.text for p in doc.paragraphs)[:10000]
            elif ext in (".xlsx", ".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(path, read_only=True)
                parts = []
                for ws in wb.worksheets[:5]:
                    parts.append(f"=== {ws.title} ===")
                    for row in ws.iter_rows(max_row=100, values_only=True):
                        parts.append("\t".join(str(c) if c is not None else "" for c in row))
                return "\n".join(parts)[:10000]
            else:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()[:10000]
        except Exception as e:
            return f"Error reading document: {e}"

    async def _handle_create_presentation(self, args: dict, **kwargs) -> str:
        """Create a PowerPoint presentation."""
        from pptx import Presentation as PptxPresentation
        from pptx.util import Pt
        output_path = args.get("output_file", os.path.join(self._project_dir, "presentation.pptx"))
        valid, resolved = _validate_file_path(output_path)
        if not valid:
            return f"Error: {resolved}"
        output_path = resolved
        slides_data = args.get("slides", [])
        title = args.get("title", "Presentation")
        try:
            prs = PptxPresentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = args.get("subtitle", "")
            for sd in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = sd.get("title", "")
                body = slide.placeholders[1]
                tf = body.text_frame
                for bullet in sd.get("bullets", []):
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            prs.save(output_path)
            size_kb = os.path.getsize(output_path) / 1024
            return json.dumps({"success": True, "path": output_path, "slides": len(slides_data) + 1, "size_kb": round(size_kb, 1)})
        except Exception as e:
            return f"Error creating presentation: {e}"

    # ── Media Tools ──────────────────────────────────────────────────────────

    async def _handle_image_edit(self, args: dict, **kwargs) -> str:
        """Edit an image: resize, crop, rotate, convert format."""
        from PIL import Image
        input_path = args.get("input_file", "")
        if not input_path:
            return "Error: input_file is required"
        valid, resolved = _validate_file_path(input_path)
        if not valid:
            return f"Error: {resolved}"
        input_path = resolved
        output_path = args.get("output_file", input_path)
        valid2, resolved2 = _validate_file_path(output_path)
        if not valid2:
            return f"Error: {resolved2}"
        output_path = resolved2
        try:
            img = Image.open(input_path)
            operation = args.get("operation", "resize")
            if operation == "resize":
                w = args.get("width", img.width)
                h = args.get("height", img.height)
                img = img.resize((int(w), int(h)), Image.LANCZOS)
            elif operation == "crop":
                left = args.get("left", 0)
                top = args.get("top", 0)
                right = args.get("right", img.width)
                bottom = args.get("bottom", img.height)
                img = img.crop((left, top, right, bottom))
            elif operation == "rotate":
                angle = args.get("angle", 90)
                img = img.rotate(angle, expand=True)
            elif operation == "convert":
                pass
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            img.save(output_path)
            return json.dumps({"success": True, "path": output_path, "size": f"{img.width}x{img.height}"})
        except Exception as e:
            return f"Error editing image: {e}"

    async def _handle_render_diagram(self, args: dict, **kwargs) -> str:
        """Render a diagram file (.mmd, .d2, .puml) to PNG."""
        input_path = args.get("input_file", "")
        output_path = args.get("output_file", "")
        if not input_path:
            return "Error: input_file is required"
        valid, resolved = _validate_file_path(input_path)
        if not valid:
            return f"Error: {resolved}"
        input_path = resolved
        if not output_path:
            output_path = input_path.rsplit(".", 1)[0] + ".png"
        valid2, resolved2 = _validate_file_path(output_path)
        if not valid2:
            return f"Error: {resolved2}"
        output_path = resolved2
        try:
            ext = os.path.splitext(input_path)[1].lower()
            if ext == ".mmd":
                cmd = f"mmdc -i '{input_path}' -o '{output_path}' -b transparent -p /root/puppeteer-config.json"
            elif ext == ".puml":
                cmd = f"plantuml -tpng '{input_path}' -o '{os.path.dirname(output_path)}'"
            elif ext == ".d2":
                cmd = f"d2 '{input_path}' '{output_path}'"
            else:
                cmd = f"mmdc -i '{input_path}' -o '{output_path}' -p /root/puppeteer-config.json"
            result = await execute_sandboxed(command=cmd, working_dir=os.path.dirname(input_path), timeout=30)
            if os.path.exists(output_path):
                return json.dumps({"success": True, "path": output_path})
            return f"Error rendering diagram: {result.get('stderr', 'Unknown error')}"
        except Exception as e:
            return f"Error: {e}"

    async def _handle_generate_chart(self, args: dict, **kwargs) -> str:
        """Generate a chart using matplotlib and save as image."""
        chart_type = args.get("chart_type", "bar")
        data = args.get("data", {})
        title = args.get("title", "Chart")
        output_path = args.get("output_file", os.path.join(self._project_dir, "chart.png"))
        valid, resolved = _validate_file_path(output_path)
        if not valid:
            return f"Error: {resolved}"
        output_path = resolved
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            fig, ax = plt.subplots(figsize=(10, 6))
            labels = data.get("labels", [])
            values = data.get("values", [])
            if chart_type == "bar":
                ax.bar(labels, values, color=data.get("colors", None))
            elif chart_type == "line":
                ax.plot(labels, values, marker="o")
            elif chart_type == "pie":
                ax.pie(values, labels=labels, autopct="%1.1f%%")
            elif chart_type == "scatter":
                x = data.get("x", labels)
                y = data.get("y", values)
                ax.scatter(x, y)
            ax.set_title(title)
            if chart_type != "pie":
                ax.set_xlabel(data.get("xlabel", ""))
                ax.set_ylabel(data.get("ylabel", ""))
            plt.tight_layout()
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
            plt.savefig(output_path, dpi=150, bbox_inches="tight")
            plt.close()
            return json.dumps({"success": True, "path": output_path})
        except Exception as e:
            return f"Error generating chart: {e}"

    # ── Advanced Shell Tools ─────────────────────────────────────────────────

    async def _handle_shell_send(self, args: dict, **kwargs) -> str:
        """Send input to a running shell process (stdin)."""
        session = self._get_session(args.get("session", "default"))
        input_text = args.get("input", "")
        command = args.get("command", "")
        if command:
            result = await session.execute(f"echo '{input_text}' | {command}", timeout=args.get("timeout", 10))
            return result.get("stdout", "") + (f"\nSTDERR: {result['stderr']}" if result.get("stderr") else "")
        return "Error: command is required for shell_send"

    async def _handle_shell_wait(self, args: dict, **kwargs) -> str:
        """Wait for a command to complete (check last session result)."""
        session = self._get_session(args.get("session", "default"))
        if session._history:
            last = session._history[-1]
            return f"Last command: {last['command']}\nExit code: {last['exit_code']}\nOutput: {last['stdout'][:3000]}"
        return "No commands in session history."

    async def _handle_shell_kill(self, args: dict, **kwargs) -> str:
        """Kill a process by PID or pattern."""
        pid = args.get("pid")
        pattern = args.get("pattern", "")
        if pid:
            cmd = f"kill -9 {pid}"
        elif pattern:
            cmd = f"pkill -f '{pattern}'"
        else:
            return "Error: pid or pattern is required"
        session = self._get_session("kill")
        result = await session.execute(cmd, timeout=5)
        return f"Kill result: exit_code={result['exit_code']}\n{result.get('stdout', '')}{result.get('stderr', '')}"

    # ── Advanced Browser Tools ───────────────────────────────────────────────

    async def _handle_browser_fill_form(self, args: dict, **kwargs) -> str:
        """Fill multiple form fields at once."""
        browser = await self._get_browser()
        fields = args.get("fields", [])
        if not fields:
            return "Error: fields array is required"
        results = []
        for field in fields:
            index = field.get("index")
            value = field.get("value", "")
            try:
                result = await browser.input_text(value, index=index)
                results.append(f"Field {index}: OK")
            except Exception as e:
                results.append(f"Field {index}: Error - {e}")
        return "\n".join(results)

    async def _handle_browser_move_mouse(self, args: dict, **kwargs) -> str:
        """Move mouse cursor to a position (for hover effects)."""
        browser = await self._get_browser()
        x = args.get("coordinate_x", 0)
        y = args.get("coordinate_y", 0)
        await browser.initialize()
        try:
            await browser._page.mouse.move(float(x), float(y))
            await asyncio.sleep(0.3)
            screenshot = await browser._take_screenshot()
            return f"Mouse moved to ({x}, {y}). Screenshot: {screenshot}"
        except Exception as e:
            return f"Mouse move error: {e}"

    async def _handle_browser_close(self, args: dict, **kwargs) -> str:
        """Close the browser session."""
        if self._browser:
            try:
                await self._browser.close()
            except Exception:
                pass
            self._browser = None
        return "Browser closed."

    # ── Infrastructure Tools ─────────────────────────────────────────────────

    async def _handle_expose_port(self, args: dict, **kwargs) -> str:
        """Expose a local port for temporary public access."""
        port = args.get("port", 8080)
        try:
            result = await execute_sandboxed(
                command=f"which cloudflared && cloudflared tunnel --url http://localhost:{port} &",
                working_dir="/tmp",
                timeout=10,
            )
            if result["exit_code"] == 0:
                return json.dumps({"success": True, "port": port, "method": "cloudflared", "note": "Tunnel starting in background"})
            return json.dumps({
                "success": True,
                "port": port,
                "url": f"http://localhost:{port}",
                "note": "Port is accessible locally. Use SSH tunnel or reverse proxy for external access."
            })
        except Exception as e:
            return f"Error exposing port: {e}"

    async def _handle_speech_to_text(self, args: dict, **kwargs) -> str:
        """Transcribe audio/video to text."""
        input_path = args.get("input_file", "")
        if not input_path:
            return "Error: input_file is required"
        valid, resolved = _validate_file_path(input_path)
        if not valid:
            return f"Error: {resolved}"
        input_path = resolved
        try:
            result = await execute_sandboxed(
                command=f'python3 -c "import whisper; m=whisper.load_model(\'base\'); r=m.transcribe(\'{input_path}\'); print(r[\'text\'])"',
                working_dir=os.path.dirname(input_path),
                timeout=120,
            )
            if result["exit_code"] == 0 and result["stdout"].strip():
                return result["stdout"].strip()
            return f"Transcription not available. Install whisper: pip install openai-whisper. Error: {result.get('stderr', '')[:500]}"
        except Exception as e:
            return f"Error: {e}"

    async def _handle_text_to_speech(self, args: dict, **kwargs) -> str:
        """Convert text to speech audio file."""
        text = args.get("text", "")
        output_path = args.get("output_file", os.path.join(self._project_dir, "speech.mp3"))
        if not text:
            return "Error: text is required"
        valid, resolved = _validate_file_path(output_path)
        if not valid:
            return f"Error: {resolved}"
        output_path = resolved
        try:
            import httpx
            api_key = os.environ.get("OPENAI_API_KEY", "")
            if api_key:
                async with httpx.AsyncClient() as client:
                    resp = await client.post(
                        "https://api.openai.com/v1/audio/speech",
                        headers={"Authorization": f"Bearer {api_key}"},
                        json={"model": "tts-1", "input": text[:4096], "voice": args.get("voice", "alloy")},
                        timeout=60,
                    )
                    if resp.status_code == 200:
                        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
                        with open(output_path, "wb") as f:
                            f.write(resp.content)
                        return json.dumps({"success": True, "path": output_path, "size_kb": round(len(resp.content)/1024, 1)})
            return "Error: TTS requires OPENAI_API_KEY"
        except Exception as e:
            return f"Error: {e}"

    async def _handle_parallel_map(self, args: dict, **kwargs) -> str:
        """Execute parallel subtasks (simplified — runs sequentially for safety)."""
        inputs = args.get("inputs", [])
        if not inputs:
            return "Error: inputs array is required"
        return json.dumps({
            "status": "parallel_map_not_yet_implemented",
            "inputs_count": len(inputs),
            "note": "Use shell_exec with background processes (&) for parallel execution, or iterate sequentially."
        })

    # ═══════════════════════════════════════════════════════════════════════════
    # SKILLS SYSTEM — modular knowledge base
    # ═══════════════════════════════════════════════════════════════════════════

    _SKILLS_DIR = "/root/arcane/skills"

    async def _handle_read_skill(self, args: dict, **kwargs) -> str:
        """Read a skill's SKILL.md file for best practices and instructions."""
        skill_name = args.get("name", "")
        if not skill_name:
            return "Error: skill name is required"
        # Sanitize
        skill_name = skill_name.replace("/", "").replace("..", "")
        skill_path = os.path.join(self._SKILLS_DIR, skill_name, "SKILL.md")
        if not os.path.exists(skill_path):
            # Try fuzzy match
            available = []
            if os.path.isdir(self._SKILLS_DIR):
                available = [d for d in os.listdir(self._SKILLS_DIR)
                           if os.path.isdir(os.path.join(self._SKILLS_DIR, d))]
            return json.dumps({
                "error": f"Skill '{skill_name}' not found",
                "available_skills": available,
            })
        try:
            with open(skill_path, "r", encoding="utf-8") as f:
                content = f.read()
            return content
        except Exception as e:
            return f"Error reading skill: {e}"

    async def _handle_list_skills(self, args: dict, **kwargs) -> str:
        """List all available skills."""
        if not os.path.isdir(self._SKILLS_DIR):
            return json.dumps({"skills": [], "note": "Skills directory not found"})
        skills = []
        for name in sorted(os.listdir(self._SKILLS_DIR)):
            skill_dir = os.path.join(self._SKILLS_DIR, name)
            if os.path.isdir(skill_dir):
                skill_md = os.path.join(skill_dir, "SKILL.md")
                description = ""
                if os.path.exists(skill_md):
                    with open(skill_md, "r") as f:
                        # Extract first heading as description
                        for line in f:
                            if line.startswith("# "):
                                description = line.strip("# \n")
                                break
                skills.append({"name": name, "description": description})
        return json.dumps({"skills": skills, "count": len(skills)})

    # ═══════════════════════════════════════════════════════════════════════════
    # WEBDEV SCAFFOLDING
    # ═══════════════════════════════════════════════════════════════════════════

    _SCAFFOLDS = {
        "react-vite": {
            "commands": [
                "npm create vite@latest {name} -- --template react-ts",
                "cd {name} && npm install",
                "cd {name} && npm install -D tailwindcss @tailwindcss/vite",
            ],
            "description": "Vite + React + TypeScript + TailwindCSS",
        },
        "static": {
            "commands": [
                "mkdir -p {name}/assets/images {name}/assets/fonts {name}/css {name}/js",
                "echo '<!DOCTYPE html><html lang=\"en\"><head><meta charset=\"UTF-8\"><meta name=\"viewport\" content=\"width=device-width,initial-scale=1.0\"><title>{name}</title><link rel=\"stylesheet\" href=\"css/style.css\"></head><body><h1>{name}</h1><script src=\"js/main.js\"></script></body></html>' > {name}/index.html",
                "echo '/* Main styles */' > {name}/css/style.css",
                "echo '// Main JS' > {name}/js/main.js",
            ],
            "description": "Plain HTML/CSS/JS",
        },
        "fastapi": {
            "commands": [
                "mkdir -p {name}/app/routers {name}/app/models {name}/app/schemas {name}/migrations",
                "echo 'fastapi\\nuvicorn\\nsqlalchemy\\nalembic\\npydantic' > {name}/requirements.txt",
                "cd {name} && python3 -m venv venv && . venv/bin/activate && pip install -r requirements.txt 2>/dev/null",
            ],
            "description": "FastAPI + SQLAlchemy backend",
        },
    }

    async def _handle_init_project(self, args: dict, **kwargs) -> str:
        """Initialize a new web project with scaffolding."""
        name = args.get("name", "my-app")
        scaffold = args.get("scaffold", "static")
        description = args.get("description", "")
        # Sanitize name
        name = name.replace("/", "").replace("..", "").replace(" ", "-")
        if scaffold not in self._SCAFFOLDS:
            return json.dumps({
                "error": f"Unknown scaffold: {scaffold}",
                "available": list(self._SCAFFOLDS.keys()),
            })
        project_dir = os.path.join(self._project_dir, name)
        if os.path.exists(project_dir):
            return json.dumps({"error": f"Directory already exists: {project_dir}"})
        scaffold_info = self._SCAFFOLDS[scaffold]
        session = self._get_session("scaffold")
        results = []
        for cmd_template in scaffold_info["commands"]:
            cmd = cmd_template.format(name=name)
            result = await session.execute(cmd, timeout=60, working_dir=self._project_dir)
            results.append({
                "command": cmd[:100],
                "exit_code": result["exit_code"],
                "output": result.get("stdout", "")[:200],
            })
            if result["exit_code"] != 0:
                return json.dumps({
                    "error": f"Scaffold step failed: {cmd[:100]}",
                    "stderr": result.get("stderr", "")[:500],
                    "results": results,
                })
        return json.dumps({
            "success": True,
            "project_dir": project_dir,
            "scaffold": scaffold,
            "description": scaffold_info["description"],
            "steps_completed": len(results),
        })

    # ═══════════════════════════════════════════════════════════════════════════
    # ENHANCED SEARCH (multi-type)
    # ═══════════════════════════════════════════════════════════════════════════

    async def _handle_search(self, args: dict, **kwargs) -> str:
        """Enhanced search with type-specific behavior."""
        query = args.get("query", "")
        search_type = args.get("type", "info")
        if not query:
            return "Error: query is required"
        # Delegate to web_search with type annotation
        # The underlying web_search handler already exists — enhance it with type context
        enhanced_query = query
        if search_type == "api":
            enhanced_query = f"{query} API documentation"
        elif search_type == "news":
            enhanced_query = f"{query} latest news"
        elif search_type == "research":
            enhanced_query = f"{query} research paper"
        elif search_type == "data":
            enhanced_query = f"{query} dataset download"
        elif search_type == "image":
            enhanced_query = f"{query} high quality images"
        # Use existing web_search handler
        result = await self._handle_web_search({"query": enhanced_query})
        return result
