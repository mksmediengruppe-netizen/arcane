from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Цвета и размеры шрифтов
COLOR_TITLE = RGBColor(26, 26, 46)
COLOR_TEXT = RGBColor(33, 37, 41)
COLOR_SUBTLE = RGBColor(102, 112, 128)

FS_TITLE = Pt(44)
FS_SUBTITLE = Pt(22)
FS_SECTION = Pt(36)
FS_BULLET = Pt(24)


def style_title(shape, size=FS_SECTION):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT


def style_paragraph(p, size=FS_BULLET, color=COLOR_TEXT, level=0):
    p.level = level
    p.font.size = size
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = title
    style_title(slide.shapes.title, size=FS_TITLE)

    sub = slide.placeholders[1]
    sub.text = subtitle
    p = sub.text_frame.paragraphs[0]
    p.font.size = FS_SUBTITLE
    p.font.color.rgb = COLOR_SUBTLE
    p.alignment = PP_ALIGN.LEFT


def add_bullets_slide(prs: Presentation, title: str, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    style_title(slide.shapes.title, size=FS_SECTION)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, text in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        style_paragraph(p, size=FS_BULLET, color=COLOR_TEXT, level=0)


def build_presentation(output_path: str):
    prs = Presentation()

    # 1) Титульный
    add_title_slide(
        prs,
        title="Преимущества облачных технологий",
        subtitle=(
            "Краткая презентация: что такое облако, ключевые преимущества и примеры использования\n"
            "Подготовлено ARCANE"
        ),
    )

    # 2) Что такое облако
    add_bullets_slide(
        prs,
        title="Что такое облако",
        bullets=[
            "Модель предоставления вычислительных ресурсов через интернет по запросу",
            "Модели сервисов: IaaS, PaaS, SaaS",
            "Ключевые свойства: эластичность, самообслуживание, оплата по факту потребления",
        ],
    )

    # 3) Преимущества
    add_bullets_slide(
        prs,
        title="Преимущества",
        bullets=[
            "Масштабируемость под пиковые нагрузки",
            "CapEx → OpEx: снижение капитальных затрат",
            "Высокая доступность и отказоустойчивость",
            "Глобальная инфраструктура: регионы, CDN",
            "Безопасность и соответствие стандартам",
            "Автоматизация и поддержка DevOps-практик",
        ],
    )

    # 4) Примеры использования
    add_bullets_slide(
        prs,
        title="Примеры использования",
        bullets=[
            "Резервное копирование и аварийное восстановление (DR)",
            "Веб‑приложения и микросервисы (Kubernetes, serverless)",
            "Аналитика и Big Data (Data Lake, ML‑пайплайны)",
            "IoT и обработка событий в реальном времени",
            "CI/CD и изолированные тестовые окружения",
        ],
    )

    # 5) Заключение
    add_bullets_slide(
        prs,
        title="Заключение",
        bullets=[
            "Облако ускоряет инновации и снижает TCO",
            "Начинайте с пилота и принципов Well‑Architected",
            "Выбирайте провайдера под требования: AWS, Azure, GCP",
            "Вопросы?",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_presentation("/root/workspace/0b9cdb61-7f71-43f8-9568-e27d8b0944e9/cloud_tech_benefits.pptx")
